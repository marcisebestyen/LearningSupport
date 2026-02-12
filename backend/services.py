import io
import re
import fitz
import google.generativeai as genai
from gtts import gTTS
from sqlalchemy.orm import Session, joinedload
import models
import os
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from google.auth.transport.requests import Request
import json
from google.generativeai.types import GenerationConfig
from docx import Document as DocxDocument
from pptx import Presentation
from sqlalchemy import text as sql_text

genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# --- Document services ---

class DocumentService:
    def __init__(self):
        genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
        self.model = genai.GenerativeModel('gemini-2.5-flash')
        self.embedding_model = "models/gemini-embedding-001"

    def extract_text(self, file_bytes: bytes, filename: str) -> str:
        ext = filename.split('.')[-1].lower()
        if ext == 'pdf':
            return self._extract_from_pdf(file_bytes)
        elif ext == 'docx':
            return self._extract_from_docx(file_bytes)
        elif ext == 'pptx':
            return self._extract_from_pptx(file_bytes)
        else:
            raise ValueError(f"Unsupported file type: {ext}")


    def _extract_from_pdf(self, file_bytes: bytes):
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        return "".join([page.get_text() for page in doc])


    def _extract_from_docx(self, file_bytes: bytes):
        cod = DocxDocument(io.BytesIO(file_bytes))
        return "\n".join([para.text for para in cod.paragraphs])


    def _extract_from_pptx(self, file_bytes: bytes):
        prs = Presentation(io.BytesIO(file_bytes))
        text_content = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)

        return "\n".join(text_content)


    def generate_summary(self, text: str, references: list[str] = None, study_focus: str = None) -> str:
        ref_text = ""
        if references:
            ref_text = "\n\n### 📚 Ajálott irodalom & Források\n" + "\n".join([f"* {ref}" for ref in references])

        if study_focus:
            instruction = f"""
            The user has uploaded a large document but ONLY wants to focus on this topic: "{study_focus}".
            
            TASK:
            1. Scan the text provided.
            2. Locate the sections relevant to "{study_focus}".
            3. Ignore unrelated chapters/sections.
            4. Create a summary ONLY for the requested topic in HUNGARIAN.
            """
        else:
            instruction = "Analyze the following document. Output the summary explicitly in HUNGARIAN."

        prompt = f"""
        {instruction}
        Structure it for a student.
        Text content (partial): {text[:50000]}
        
        IMPORTANT FORMATTING RULES:
        1. Do NOT use code blocks (```) for the text.
        2. Use standard bullet points (*).
        3. Do NOT indent nested bullet points with more than 2 spaces.
        4. Use bolding (**text**) for key terms.
        """
        try:
            response = self.model.generate_content(prompt)
            summary = response.text
            return summary + ref_text
        except Exception as e:
            print(f"AI Error: {e}")
            return "Hiba történt az összefoglaló generálása közben."


    def save_document(self, db: Session, filename: str, content: str, summary: str, user_id: int, category: str = None, study_focus: str = None):
        new_doc = models.Document(
            filename=filename,
            content=content,
            summary=summary,
            owner_id=user_id,
            category=category,
            study_focus=study_focus
        )
        db.add(new_doc)
        db.commit()
        db.refresh(new_doc)

        chunks = self._chunk_text(content)
        for idx, chunk_text in enumerate(chunks):
            vector = self._get_embedding(chunk_text)

            db_chunk = models.DocumentChunk(
                document_id=new_doc.id,
                chunk_index=idx,
                content=chunk_text,
                embedding=vector,
            )
            db.add(db_chunk)

        db.commit()
        return new_doc


    def _chunk_text(self, text: str, chunk_size: int = 1000) -> list[str]:
        return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]


    def _get_embedding(self, text: str):
        result = genai.embed_content(
            model=self.embedding_model,
            content=text,
            # task_type="retrieval_document"
            output_dimensionality=768
        )
        return result['embedding']


    def get_user_history(self, db: Session, user_id: int):
        return db.query(models.Document).filter(models.Document.owner_id == user_id).all()


    def delete_document(self, db: Session, doc_id: int, user_id: int):
        doc = db.query(models.Document).filter(
            models.Document.id == doc_id,
            models.Document.owner_id == user_id
        ).first()

        if doc:
            db.delete(doc)
            db.commit()
            return True

        return False


    def process_audio_generation(self, db: Session, doc: models.Document):
        try:
            clean_text = re.sub(r'[*_#]', '', doc.summary)
            clean_text = clean_text.replace("\n", " ")

            tts = gTTS(text=clean_text, lang='hu')

            temp_path = f"temp_audio_{doc.id}.mp3"
            tts.save(temp_path)

            drive_service = GoogleDriveService()
            drive_id = drive_service.upload_audio(temp_path, f"Summary_{doc.filename}.mp3")

            doc.google_drive_id = drive_id
            db.commit()

            if os.path.exists(temp_path):
                os.remove(temp_path)

            return drive_id
        except Exception as e:
            print(f"gTTS Error: {e}")
            raise e


    def validate_content(self, text: str) -> dict:
        prompt = f"""
                Act as a strict Fact-Checker and Librarian. Analyze the text below (first 15000 chars).

                Tasks:
                1. VALIDATION: Determine if this is a coherent, factually possible text (scientific, fictional story, history, etc.) OR if it is incoherent/blatant generated nonsense/fake news.
                2. RESOURCES: If it is Valid, provide 3 real-world book titles or reliable URLs relevant to the topic.

                Text Preview:
                {text[:15000]}

                Output JSON format ONLY:
                {{
                    "is_valid": boolean,
                    "warning_message": "String explaining why it looks fake (or null if valid)",
                    "references": ["Title/URL 1", "Title/URL 2", "Title/URL 3"]
                }}
                """

        try:
            config = GenerationConfig(response_mime_type="application/json")
            response = self.model.generate_content(prompt, generation_config=config)
            return json.loads(response.text)
        except Exception as e:
            print("Validation error: ", e)
            return {
                "is_valid": True,
                "warning_message": "⚠️ Technical Error: Validation could not be completed. The document content was not verified.",
                "references": []
            }


    def find_cross_references(self, db: Session, current_doc_id: int, content: str, user_id: int):
        new_embedding = self._get_embedding(content[:2000])

        query = sql_text(f"""
        SELECT d.filename, dc.content, d.id
        FROM document_chunks dc
            JOIN documents d ON d.id = dc.document_id
        WHERE d.id != :current_doc_id
            AND d.owner_id = :user_id
        ORDER BY d.embedding <=> :embedding
        LIMIT 3 
        """)

        results = db.execute(query, {
            "user_id": user_id,
            "current_doc_id": current_doc_id,
            "embedding": str(new_embedding)
        }).fetchall()
        if not results:
            return ""

        cross_context = "\n".join([f"Source ({r[0]}): {r[1]}" for r in results])

        prompt = f"""
                You are a 'Second Brain' assistant.
                Analyze connections between the new document and previous ones.

                New Document: {content[:3000]}
                Previous Snippets: {cross_context}

                TASK:
                Write a concise 'Cross-Reference' note (Hungarian).
                1. Explain the connection/similarity.
                2. Reference the old filenames explicitly.
                3. Do NOT include a Title or Header. Start directly with the text.
                """

        try:
            response = self.model.generate_content(prompt)
            ai_text = response.text.strip()

            ai_text_html = ai_text.replace("\n", "<br>")

            ai_text_html = re.sub(r'^#+\s*', '', ai_text_html)
            ai_text_html = re.sub(r'\*\*', '', ai_text_html)

            formatted_html = (
                f'<div class="second-brain-container">'
                f'<div class="sb-header">'
                f'<span class="sb-icon">🧠</span>'
                f'<span class="sb-title">Másodlagos Agy &middot; Kapcsolódó Ismeretek</span>'
                f'</div>'
                f'<div class="sb-content">{ai_text_html}</div>'
                f'</div>'
            )

            return formatted_html

        except Exception as e:
            print(f"Error getting cross reference: {e}")
            return ""


# --- User services

class UserService:
    def create_user(self, db: Session, username: str, hashed_pass: str):
        user = models.User(username=username, hashed_password=hashed_pass)
        db.add(user)
        db.commit()
        db.refresh(user)
        return user


    def get_user_by_username(self, db: Session, username: str):
        return db.query(models.User).filter(models.User.username == username).first()

# --- Quiz and Flashcard services

class QuizService:
    def __init__(self, db: Session = None):
        self.db = db
        self.model = genai.GenerativeModel('gemini-2.5-flash')


    def generate_quiz(self, document_id: int, user_id: int):
        doc: models.Document | None = self.db.query(models.Document).filter(
            models.Document.id == document_id,
            models.Document.owner_id == user_id
        ).first()

        if not doc:
            print(f"Error: Document {document_id} not found.")
            return None

        prompt = f"""
        Generate a quiz based on the text below.
        Language: HUNGARIAN.
        Format: JSON Array of objects.
        Quantity: 10 Questions.
        Types: Mix of "multiple_choice" and "true_false".

        Schema for each object:
        {{
            "question_text": "The question string",
            "type": "multiple_choice" OR "true_false",
            "options": ["Option A", "Option B", "Option C", "Option D"] (OR null if true_false),
            "correct_answer": "The exact string matching one of the options or 'Igaz'/'Hamis'"
        }}

        Ensure the "correct_answer" exactly matches one of the strings in "options".

        Text content:
        {doc.content[:30000]}
        """

        try:
            config = GenerationConfig(response_mime_type="application/json")
            response = self.model.generate_content(
                prompt,
                generation_config=config
            )

            raw_text = response.text.strip()
            if raw_text.startswith("```"):
                raw_text = raw_text.split("\n", 1)[1]
                if raw_text.endswith("```"):
                    raw_text = raw_text.rsplit("\n", 1)[0]

            quiz_data = json.loads(raw_text)

            if isinstance(quiz_data, dict):
                if "questions" in quiz_data:
                    quiz_data = quiz_data["questions"]
                elif "quiz" in quiz_data:
                    quiz_data = quiz_data["quiz"]
                else:
                    for val in quiz_data.values():
                        if isinstance(val, list):
                            quiz_data = val
                            break

            new_quiz = models.Quiz(
                document_id=doc.id,
                owner_id=user_id,
                top_score=0,
                passed=False
            )
            self.db.add(new_quiz)
            self.db.commit()
            self.db.refresh(new_quiz)

            for q_data in quiz_data:
                options = q_data.get('options')
                if q_data.get('type') == 'true_false' and not options:
                    options = ["Igaz", "Hamis"]

                question = models.Question(
                    quiz_id=new_quiz.id,
                    question_text=q_data['question_text'],
                    question_type=q_data['type'],
                    options=options,
                    correct_answer=q_data['correct_answer']
                )
                self.db.add(question)

            self.db.commit()
            print(f"--- Quiz Generated ID: {new_quiz.id} ---")
            return new_quiz

        except Exception as e:
            print(f"!!! Quiz Generation Error: {e}")
            return None


    def generate_flashcards(self, document_id: int, user_id: int):
        doc = self.db.query(models.Document).filter(
            models.Document.id == document_id,
            models.Document.owner_id == user_id
        ).first()
        if not doc:
            return None

        prompt = f"""
        Analyze the text below and generate 10 flashcards. 
        Output PURE JSON: [{{ "front": "term", "back": "definition" }}]
        Language: HUNGARIAN.
        Text: {doc.content[:30000]}
        """

        try:
            config = GenerationConfig(response_mime_type="application/json")
            response = self.model.generate_content(prompt, generation_config=config)
            cards_data = json.loads(response.text)

            new_set = models.FlashcardSet(document_id=doc.id)
            self.db.add(new_set)
            self.db.commit()
            self.db.refresh(new_set)

            for card in cards_data:
                db_card = models.Flashcard(set_id=new_set.id, front=card['front'], back=card['back'])
                self.db.add(db_card)

            self.db.commit()
            return new_set.id

        except Exception as e:
            print(f"!!! Flashcard Generation Error: {e}")
            return None


    def get_flashcard_set(self, set_id: int):
        return self.db.query(models.FlashcardSet).options(
            joinedload(models.FlashcardSet.cards)
        ).filter(models.FlashcardSet.id == set_id).first()


    def get_user_flashcard_sets(self, user_id: int):
        sets = self.db.query(models.FlashcardSet).join(models.Document).filter(
            models.Document.owner_id == user_id
        ).order_by(models.FlashcardSet.created_at.desc()).all()

        results = []
        for fset in sets:
            results.append({
                "id": fset.id,
                "created_at": fset.created_at,
                "document_filename": fset.document.filename,
                "card_count": len(fset.cards),
            })
        return results


    def get_user_quizzes(self, user_id: int):
        quizzes = self.db.query(models.Quiz).options(
            joinedload(models.Quiz.document)
        ).filter(
            models.Quiz.owner_id == user_id
        ).order_by(models.Quiz.created_at.desc()).all()

        return [q for q in quizzes if q.document is not None]


    def get_quiz_by_id(self, quiz_id: int, user_id: int):
        return self.db.query(models.Quiz).options(
            joinedload(models.Quiz.questions)
        ).filter(
            models.Quiz.id == quiz_id,
            models.Quiz.owner_id == user_id
        ).first()


    def submit_score(self, quiz_id: int, score: int, user_id: int):
        quiz = self.db.query(models.Quiz).filter(
            models.Quiz.id == quiz_id,
            models.Quiz.owner_id == user_id
        ).first()

        if quiz:
            if score > quiz.top_score:
                quiz.top_score = score

            if score >= 5:
                quiz.passed = True

            self.db.commit()
            self.db.refresh(quiz)
            return quiz
        return None

# --- Chat services

class ChatService:
    def __init__(self, db: Session):
        self.db = db
        self.model = genai.GenerativeModel('gemini-2.5-flash')

    def ask_document(self, doc_id: int, question: str):
        usr_msg = models.ChatMessage(document_id=doc_id, role='user', content=question)
        self.db.add(usr_msg)
        self.db.commit()

        q_embedding = genai.embed_content(
            model='models/gemini-embedding-001',
            content=question,
            # task_type='retrieval_query'
            output_dimensionality=768
        )['embedding']

        query = sql_text("""
            SELECT content
            FROM document_chunks
            WHERE document_id = :doc_id
            ORDER BY embedding <=> :q_embedding
            LIMIT 3
        """)

        results = self.db.execute(query, {"doc_id": doc_id, "q_embedding": str(q_embedding)}).fetchall()
        context_text = "\n\n".join([row[0] for row in results])

        prompt = f"""
        You are a helpful tutor. Answer the question based ONLY on the context below.
        
        Context: 
        {context_text}
        
        Question: {question}
        Answer (in Hungarian):
        """

        response = self.model.generate_content(prompt)
        answer_text = response.text

        ai_msg = models.ChatMessage(document_id=doc_id, role='ai', content=answer_text)
        self.db.add(ai_msg)
        self.db.commit()

        return answer_text


    def get_chat_history(self, doc_id: int, mode: str = 'chat'):
        query = self.db.query(models.ChatMessage).filter(
            models.ChatMessage.document_id == doc_id,
        )

        if mode == 'tutor':
            query = query.filter(models.ChatMessage.role.in_(['tutor_ai', 'tutor_user']))
        else:
            query = query.filter(models.ChatMessage.role.in_(['ai', 'user']))

        return query.order_by(models.ChatMessage.created_at.asc()).all()

    def start_socratic_session(self, doc_id: int):
        doc = self.db.query(models.Document).filter(
            models.Document.id == doc_id
        ).first()
        if not doc:
            return None

        existing_history = self.db.query(models.ChatMessage).filter(
            models.ChatMessage.document_id == doc_id,
            models.ChatMessage.role.in_(['tutor_ai', 'tutor_user'])
        ).all()
        if existing_history:
            return existing_history[-1].content

        prompt = f"""
                You are a Socratic Tutor. Your goal is to test the student's understanding of the text below.

                RULES:
                1. Do NOT summarize the text.
                2. Ask ONE specific, open-ended question based on the text to start the discussion.
                3. The question should require understanding, not just copy-pasting.
                4. Output language: HUNGARIAN.

                Text to teach:
                {doc.content[:30000]}

                Start with the first question now.
                """

        response = self.model.generate_content(prompt)
        ai_question = response.text

        ai_msg = models.ChatMessage(document_id=doc_id, role='tutor_ai', content=ai_question)
        self.db.add(ai_msg)
        self.db.commit()

        return ai_question


    def handle_tutor_response(self, doc_id: int, user_answer: str):
        usr_msg = models.ChatMessage(document_id=doc_id, role='tutor_user', content=user_answer)
        self.db.add(usr_msg)
        self.db.commit()

        history = self.get_chat_history(doc_id, mode='tutor')
        is_final_turn = len(history) >= 10

        doc = self.db.query(models.Document).filter(
            models.Document.id == doc_id,
        ).first()
        context_text = doc.content[:30000]

        conversation_history = "\n".join(
            [f"{'AI' if 'ai' in msg.role else 'Student'}: {msg.content}" for msg in history[-10:]])

        if is_final_turn:
            prompt = f"""
                        The tutoring session is over. Generate a Final Report based on the student's performance.

                        Context: {context_text}
                        History: {conversation_history}

                        Output strictly JSON:
                        {{
                            "text": "## 📊 Session Report Card\\n\\n**Grade:** 8/10\\n\\n**Feedback:** ...",
                            "status": "neutral",
                            "is_finish": true
                        }}
                        Language: HUNGARIAN.
                        """
        else:
            prompt = f"""
                        You are a Socratic Tutor. Analyze the user's answer.

                        Context: {context_text}
                        History: {conversation_history}
                        User Answer: {user_answer}

                        Output strictly JSON:
                        {{
                            "status": "correct" OR "incorrect" OR "neutral",
                            "text": "Your feedback here... + Next Question",
                            "is_finish": false
                        }}
                        Language: HUNGARIAN.
                        """

        try:
            config = GenerationConfig(response_mime_type="application/json")
            response = self.model.generate_content(prompt, generation_config=config)
            response_data = json.loads(response.text)

            ai_msg = models.ChatMessage(document_id=doc_id, role='tutor_ai', content=response.text)
            self.db.add(ai_msg)
            self.db.commit()

            return response_data

        except Exception as e:
            print(f"Tutor Error: {e}")
            # Fallback
            return {"status": "neutral", "text": "Hiba történt. Folytassuk...", "is_finish": False}


    def reset_tutor_history(self, doc_id: int):
        self.db.query(models.ChatMessage).filter(
            models.ChatMessage.document_id == doc_id,
            models.ChatMessage.role.in_(['tutor_ai', 'tutor_user'])
        ).delete(synchronize_session=False)

        self.db.commit()
        return True

# Mind Map services

class MindMapService:
    def __init__(self, db: Session):
        self.db = db
        self.model = genai.GenerativeModel('gemini-2.5-flash')


    def generate_mindmap(self, doc_id: int, user_id: int):
        doc = self.db.query(models.Document).filter(
            models.Document.id == doc_id,
            models.Document.owner_id == user_id
        ).first()

        if not doc:
            return None

        prompt = f"""
                Create a hierarchical mind map using Mermaid.js `graph TD` syntax.

                STRICT RULES:
                1. Use `graph TD` (Top-Down) layout.
                2. **Do NOT** use the `mindmap` keyword.
                3. **Shapes**:
                   - Use Double Circle for the Root Node: `A((Main Topic))`
                   - Use Rounded Edges for Branches: `B(Sub Topic)`
                   - Use Stadium Shape for Leaves: `C([Detail])`
                4. **Connections**: Use standard arrows `-->`.
                5. **Sanitization**: Remove all special characters `( ) [ ] " '` from labels.
                6. Language: HUNGARIAN.

                Text to analyze:
                {doc.content[:30000]}
                """

        try:
            response = self.model.generate_content(prompt)
            script = response.text.strip()

            if script.startswith("```mermaid"):
                script = script.replace("```mermaid", "").replace("```", "")
            elif script.startswith("```"):
                script = script.replace("```", "")

            script = script.strip()

            if "graph LR" in script:
                script = script.replace("graph LR", "graph TD")

            if not script.startswith("graph"):
                script = "graph TD\n" + script

            new_map = models.MindMap(
                document_id=doc_id,
                mermaid_script=script
            )
            self.db.add(new_map)
            self.db.commit()
            self.db.refresh(new_map)

            return new_map

        except Exception as e:
            print(f"!!! MindMap Generation Error: {e}")
            return None


    def get_mindmap_by_doc(self, doc_id: int, user_id: int):
        return self.db.query(models.MindMap).join(models.Document).filter(
            models.MindMap.document_id == doc_id,
            models.Document.owner_id == user_id
        ).first()


    def get_user_mindmaps(self, user_id: int):
        maps = self.db.query(models.MindMap).options(
            joinedload(models.MindMap.document)
        ).join(models.Document).filter(
            models.Document.owner_id == user_id
        ).order_by(models.MindMap.created_at.desc()).all()

        results = []
        for m in maps:
            results.append({
                "id": m.id,
                "mermaid_script": m.mermaid_script,
                "created_at": m.created_at,
                "document_filename": m.document.filename if m.document.filename else "Unknown File",
                "document_id": m.document.id,
            })
        return results

# --- Google Drive services

class GoogleDriveService:
    def __init__(self):
        self.client_id = os.getenv("GOOGLE_DRIVE_CLIENT_ID")
        self.client_secret = os.getenv("GOOGLE_DRIVE_CLIENT_SECRET")
        self.refresh_token = os.getenv("GOOGLE_DRIVE_REFRESH_TOKEN")
        self.folder_id = os.getenv("GOOGLE_DRIVE_FOLDER_ID")

        self.creds = Credentials(
            token=None,
            refresh_token=self.refresh_token,
            client_id=self.client_id,
            client_secret=self.client_secret,
            token_uri="https://oauth2.googleapis.com/token"
        )

        if not self.creds.valid:
            self.creds.refresh(Request())

        self.service = build('drive', 'v3', credentials=self.creds)


    def upload_audio(self, file_path, filename):
        file_metadata = {
            'name': filename,
            'parents': [self.folder_id],
        }
        media = MediaFileUpload(file_path, mimetype='audio/mpeg')

        file = self.service.files().create(
            body=file_metadata,
            media_body=media,
            fields='id'
        ).execute()

        return file.get('id')


    def stream_audio(self, file_id):
        request = self.service.files().get_media(fileId=file_id)
        return request.execute()

# --- Essay / Grader services


class GraderService:
    def __init__(self, db: Session):
        self.db = db
        self.model = genai.GenerativeModel('gemini-2.5-flash')


    def evaluate_essay(self, doc_id: int, user_id: int, essay_text: str):
        doc = self.db.query(models.Document).filter(
            models.Document.id == doc_id
        ).first()

        if not doc:
            raise ValueError("Document not found.")

        prompt = f"""
                You are a strict academic professor. Your task is to grade a Student Essay based ONLY on the provided Source Material.

                Source Material:
                {doc.content[:30000]}

                Student Essay:
                {essay_text}

                TASK:
                1. Analyze the essay sentence by sentence (or logical segment).
                2. Compare each segment to the Source Material.
                3. Assign a status:
                   - "correct" (Green): Factually accurate according to source.
                   - "partial" (Yellow): Sort of correct, but missing context, slightly vague, or minor logical flaw.
                   - "incorrect" (Red): Factually wrong, hallucinated, or completely contradictory to the source.
                4. Provide a brief correction/comment for each segment.
                5. Give an overall score (0-100) and a general summary.

                OUTPUT JSON FORMAT ONLY:
                {{
                    "overall_score": 85,
                    "general_feedback": "A summary of the student's performance...",
                    "segments": [
                        {{
                            "segment_text": "The exact text from the student essay being graded.",
                            "status": "correct",
                            "feedback": "Correct. This aligns with the section on..."
                        }},
                        {{
                            "segment_text": "However, the date was 1999.",
                            "status": "incorrect",
                            "feedback": "Incorrect. The source states the date was 1990."
                        }}
                    ]
                }}

                Language: HUNGARIAN.
                """

        try:
            config = GenerationConfig(response_mime_type='application/json')
            response = self.model.generate_content(prompt, generation_config=config)

            result_json = json.loads(response.text)

            submission = models.EssaySubmission(
                document_id=doc_id,
                owner_id=user_id,
                essay_content=essay_text,
                feedback_json=result_json.get('segments', []),
                overall_score=result_json.get('overall_score', 0),
                general_feedback=result_json.get('general_feedback', "Nincs általános visszajelzés.")
            )

            self.db.add(submission)
            self.db.commit()
            self.db.refresh(submission)

            return submission

        except Exception as e:
            print(f"Grading Error: {e}")
            return None


    def get_user_essays(self, user_id: int):
        return (self.db.query(models.EssaySubmission).options(
            joinedload(models.EssaySubmission.document),
        ).filter(
            models.EssaySubmission.owner_id == user_id
        ).order_by(models.EssaySubmission.created_at.desc()).all())


    def get_essay_by_id(self, essay_id: int, user_id: int):
        return self.db.query(models.EssaySubmission).options(
            joinedload(models.EssaySubmission.document),
        ).filter(
            models.EssaySubmission.id == essay_id,
            models.EssaySubmission.owner_id == user_id
        ).first()

# --- Study Plan services ---

class StudyPlanService:
    def __init__(self, db: Session):
        self.db = db
        self.model = genai.GenerativeModel('gemini-2.5-flash')


    def generate_study_plan(self, doc_id: int, user_id: int):
        doc = self.db.query(models.Document).filter(
            models.Document.id == doc_id,
            models.Document.owner_id == user_id
        ).first()

        if not doc:
            return None

        focus_instruction = ""
        if doc.study_focus:
            focus_instruction = f"The student ONLY wants to learn about: '{doc.study_focus}'. Ignore other topics in the text."

        prompt = f"""
        Act as a professional educational consultant. 
        Create a structured Study Plan based on the text below.
        
        {focus_instruction}
        
        Rules:
        1. Analyze the complexity and length of the relevant text.
        2. Break it down into logical "Days" (e.g., 3 days, 5 days, or 2 weeks depending on depth).
        3. For each day, suggest specific sub-topics and activities (Reading, Quiz, Flashcards, Essay).
        4. Output PURE JSON format.
        5. Language: HUNGARIAN.
        
        JSON Structure:
        [
            {{
                "day": 1,
                "topic": "Introduction to...",
                "activities": ["Read Section 1", "Create Mindmap", "Review Flashcards"]
            }},
            ...
        ]
        
        Text:
        {doc.content[:60000]}
        """

        try:
            config = GenerationConfig(response_mime_type='application/json')
            response = self.model.generate_content(prompt, generation_config=config)
            plan_json = json.loads(response.text)

            existing_plan = self.db.query(models.StudyPlan).filter(
                models.StudyPlan.document_id == doc_id
            ).first()

            if existing_plan:
                existing_plan.plan_json = plan_json
                self.db.commit()
                return existing_plan
            else:
                new_plan = models.StudyPlan(
                    document_id=doc_id,
                    plan_json=plan_json,
                )
                self.db.add(new_plan)
                self.db.commit()
                self.db.refresh(new_plan)
                return new_plan

        except Exception as e:
            print(f"Study Plan error: {e}")
            return None


    def get_user_study_plans(self, user_id: int):
        plans = self.db.query(models.StudyPlan).join(models.Document).filter(
            models.Document.owner_id == user_id
        ).order_by(models.StudyPlan.created_at.desc()).all()

        results = []
        for p in plans:
            total_days = 0
            if isinstance(p.plan_json, list):
                total_days = len(p.plan_json)

            results.append({
                "id": p.document.id,
                "plan_id": p.id,
                "filename": p.document.filename,
                "study_focus": p.document.study_focus,
                "created_at": p.created_at,
                "total_days": total_days,
            })
        return results


    def get_plan_by_id(self, doc_id: int, user_id: int):
        return self.db.query(models.StudyPlan).filter(
            models.StudyPlan.document_id == doc_id,
            models.Document.owner_id == user_id
        ).first()
